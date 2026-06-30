"""Build Auto-ChemInstruct defense presentation using TSU template.
Strategy: reuse template slides (preserving backgrounds/images), replace text.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import copy, shutil

TEMPLATE = "/Users/akm/PROJ/Шаблон презентации стажировка.pptx"
OUTPUT = "/Users/akm/PROJ/auto-cheminstruct/presentation_autochem.pptx"

shutil.copy2(TEMPLATE, OUTPUT)
prs = Presentation(OUTPUT)

# Keep slide backgrounds, images, logos from template — just replace text
LAYOUT = prs.slide_layouts[2]  # Основной шаблон
EMPTY = prs.slide_layouts[1]  # Пустой слайд


def add_slide(title):
    s = prs.slides.add_slide(LAYOUT)
    for sh in s.shapes:
        if sh.name == "Shape 10":
            sh.text_frame.paragraphs[0].text = title
            break
    return s


def body(s, lines, top=1.6):
    tb = s.shapes.add_textbox(Inches(0.7), Inches(top), Inches(8.6), Inches(7.0 - top))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.space_after = Pt(6)


def mono(s, text, top=1.6, sz=10):
    tb = s.shapes.add_textbox(Inches(0.3), Inches(top), Inches(9.5), Inches(7.0 - top))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if not line.strip():
            continue
        p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.name = "Courier New"


def two_col(s, left, right, top=1.6):
    for x, lines in [(0.5, left), (5.0, right)]:
        tb = s.shapes.add_textbox(Inches(x), Inches(top), Inches(4.5), Inches(5.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.space_after = Pt(5)


# ── SLIDE 0: Update supervisor info ──
for sh in prs.slides[0].shapes:
    if sh.name == "Shape 179":
        tf = sh.text_frame
        tf.clear()
        for text, sz in [
            ("Бузаев Александр Александрович", 18),
            ("Лаборатория ИИ в химии и молекулярной инженерии (AIRI×ТГУ)", 14),
            ("Зав.лаб., к.х.н.", 14),
            ("", 14),
            ("Докладчик: Aayush Kumar, стажёр-исследователь", 14),
            ("Auto-ChemInstruct — NeurIPS Datasets & Benchmarks", 14),
        ]:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = Pt(sz)
        break

# ── SLIDE 1: Update title slide ──
for sh in prs.slides[1].shapes:
    if sh.name == "Shape 10":
        sh.text_frame.paragraphs[0].text = "Auto-ChemInstruct"
    if sh.name == "Shape 11":
        tf = sh.text_frame
        tf.clear()
        for text, sz in [
            ("Agent-Driven Synthesization of RLHF Data for Chemistry DSLMs", 16),
            ("", 10),
            ("Aayush Kumar", 14),
            ("TSU Laboratory of AI in Chemistry × AIRI Institute, Moscow", 13),
            ("NeurIPS Datasets & Benchmarks Track", 13),
        ]:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = Pt(sz)

# ── SLIDE 2: Problem (reuse template content slide) ──
s2 = prs.slides[2]
# Clear existing shapes except background images
for sh in list(s2.shapes):
    if sh.has_text_frame and sh.text_frame.text.strip():
        sh.text_frame.clear()
        # Repurpose Shape 10 (title area)
        if sh.name == "Shape 10":
            p = sh.text_frame.paragraphs[0]
            p.text = "Проблема"
            p.font.size = Pt(24)
        elif sh.name in ("Shape 11", "Shape 14"):
            pass  # skip footer text
        else:
            # Delete non-essential text shapes
            sp = sh._element
            sp.getparent().remove(sp)
# Add body text
body(
    s2,
    [
        "Существующие датасеты требуют дорогой ручной разметки экспертов",
        "Позитивное смещение — показывают успешные реакции, не объясняют неудачи",
        "LLM галлюцинируют — генерируют термодинамически невозможные реакции",
        "Статические пайплайны без обратной связи и самообучения",
        "",
        "Ключевая идея:",
        "Автоматизировать верификацию через вычислительную химию",
        "(RDKit + MMFF94) как детерминированный оракул физической истины",
    ],
)

# ── SLIDES 3-10: Content slides ──
slides_data = []

# Slide 3: Architecture
s = add_slide("4-Агентная Архитектура")
mono(
    s,
    "  ┌──────────────┐     ┌──────────────┐     ┌──────────────┐         \n"
    "  │  Hypothesis  │────→│ Verification │────→│ Compilation  │──→ DPO  \n"
    "  │    Agent     │     │    Agent     │PASS │    Agent     │   Pairs \n"
    "  │  (LLM, T=0.9)│     │(RDKit+MMFF94)│     │ (6-dim qual) │         \n"
    "  └──────┬───────┘     └──────┬───────┘     └──────────────┘         \n"
    "         ↑                    │ FAIL                                  \n"
    "         │             ┌──────┴───────┐                              \n"
    "         │             │ Reflection   │ ← CARL 4-step DAG            \n"
    "         │             └──────┬───────┘   10 failure categories      \n"
    "         │                    │                                       \n"
    "         └──── LearningContext ◄── Self-Bootstrapping (T: 1.0→0.3)   \n"
    "                                                                      \n"
    "  ◄═══ MAP-Elites: 2,600 cells · 5 mutations · 4 islands ═══►",
    top=1.5,
)

# Slide 4: Core Innovation
s = add_slide("Self-Bootstrapping: Физика как Оракул")
body(
    s,
    [
        "Generate → Verify → Reflect → Accumulate → Repeat",
        "",
        "В отличие от самооценки LLM (галлюцинации), мы связываем",
        "генерацию с детерминированными физическими симуляторами",
        "(RDKit + MMFF94) — реальная энергия, реальная геометрия.",
        "",
        "Неудачные реакции → обучающие сигналы.",
        "T cosine annealing 1.0 → 0.3 (exploration → exploitation).",
        "",
        'Пример: "Нуклеофильная атака блокирована стерическим',
        "эффектом трет-бутила. Переходное состояние невозможно.",
        'Решение: менее объёмный электрофил или SN1."',
    ],
)

# Slide 5: MAP-Elites + CARL
s = add_slide("MAP-Elites × CARL Reasoning")
two_col(
    s,
    [
        "MAP-Elites (GigaEvo):",
        "• 2,600 ячеек: 26×10×10",
        "• 5 mutation operators",
        "• 4 specialist islands",
        "• Migration: 3 elites / 10 gens",
        "• Deterministic RNG",
    ],
    [
        "CARL Chains (Maestro):",
        "• 4-шаговая DAG-цепочка:",
        "  1. Steric Analysis",
        "  2. Electronic Analysis",
        "  3. Thermo/Kinetic Analysis",
        "  4. Causal Synthesis",
        "• Шаги 1-3 — ПАРАЛЛЕЛЬНО",
        "• Chemistry-specific prompts",
    ],
)

# Slide 6: Implementation
s = add_slide("Реализация: 5 Фаз · 230 Тестов")
body(
    s,
    [
        "Phase I: Foundation — Redis, Hydra, problem dir         (140 tests)",
        "Phase II: DAG Engine — Async pipeline, Kahn sort         (163 tests)",
        "Phase III: MAP-Elites — 2,600 grid, 5 ops, 4 islands    (202 tests)",
        "Phase IV: CARL Chains — 4-step parallel DAG reflection   (219 tests)",
        "Phase V: Ablation + Paper — 7-variant, NeurIPS LaTeX    (230 tests)",
        "",
        "Стек: Python 3.13 · Pydantic v2 · RDKit+MMFF94 · LangChain",
        "Fireworks AI (MiniMax-M3) · TF-IDF+NetworkX RAG · SQLite · Docker",
        "7,530 строк кода · 2,608 строк тестов · 19 Agent Skills",
        "github.com/aayushkrm/auto-cheminstruct",
    ],
)

# Slide 7: Dataset
s = add_slide("Датасет: 172 DPO Pairs · 19 Типов Реакций")
mono(
    s,
    "                    v1.0 (DeepSeek)    v3.0 (MiniMax-M3)    Merged           \n"
    "  Pairs                  110                  62              172             \n"
    "  Pass Rate             65.9%               69.7%            67.2%            \n"
    "  Reaction Types          13                  18               19             \n"
    "  Avg Quality           0.636               0.671            0.650            \n"
    "                                                                              \n"
    "  Splits: 124 train / 6 validation / 42 test                                 \n"
    "  HuggingFace: aayushkrm/autochem-instruct                                   ",
    top=1.6,
    sz=12,
)

# Slide 8: Ablation
s = add_slide("Ablation: 7-Вариантное Исследование")
mono(
    s,
    "Variant               ME    CARL   Elites   Pass Rate   Quality   Impact      \n"
    "───────────────────────────────────────────────────────────────────────────\n"
    "Baseline               —      —      10      69.0%      0.59       —          \n"
    "CARL-Only              —      ✓      15      78.8%      0.67      +14% pass   \n"
    "MAP-Elites-Only        ✓      —      24      75.0%      0.60      2.4× elites \n"
    "Full-System            ✓      ✓      25      84.6%      0.72      3.5× +17%   \n"
    "                                                                              \n"
    "MAP-Elites масштабирует популяцию · CARL улучшает качество                   \n"
    "Full-System: 3.5× elites, +17% pass, наивысшее качество                     ",
    top=1.5,
    sz=10,
)

# Slide 9: AIRI Integration
s = add_slide("Интеграция с AIRI Frameworks")
body(
    s,
    [
        "GigaEvo (FusionBrainLab/gigaevo-core, arXiv:2511.17592):",
        "  Redis DB → redis_store.py   Async DAG → dag.py",
        "  MAP-Elites → map_elites.py   Problem I/F → problems/autochem/",
        "",
        "Maestro CARL (AIRI-Institute/maestro-core):",
        "  Event-Action-Result chains → src/carl/chain.py",
        "  Parallel DAG → Steps 1-3 execute simultaneously",
        "",
        "Собственные реализации по архитектурным паттернам AIRI.",
        "27 исследовательских документов в docs/research/.",
    ],
)

# Slide 10: Conclusion
s = add_slide("Результаты и Перспективы")
body(
    s,
    [
        "Достигнуто:",
        "",
        "✓ 4-агентный автономный пайплайн для генерации данных",
        "✓ Self-bootstrapping с физической валидацией (RDKit + MMFF94)",
        "✓ MAP-Elites: 2,600 ячеек, 5 операторов, 4 острова",
        "✓ CARL 4-шаговые цепочки причинного анализа",
        "✓ 172 DPO пар, 19 типов реакций на HuggingFace",
        "✓ 7-вариантное ablation, 230/230 тестов",
        "",
        "Будущие направления:",
        "",
        "  xTB энергетика · Распределённый MAP-Elites",
        "  Fine-tuning DSLMs · Materials science, drug discovery",
        "",
        "Auto-ChemInstruct: автономная, физически-обоснованная",
        "генерация данных возможна в масштабе.",
    ],
)

# ── Update thank-you slide (last) ──
# Add a thank-you slide using empty layout with branding
s = prs.slides.add_slide(EMPTY)

# Copy background shape from slide 0
for sh in prs.slides[0].shapes:
    if sh.name in ("Рисунок 8", "Picture 2", "Рисунок 9"):
        el = sh._element
        new_el = copy.deepcopy(el)
        s.shapes._spTree.append(new_el)

tb = s.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7.0), Inches(3.5))
tf = tb.text_frame
tf.word_wrap = True
for text, sz in [
    ("Спасибо за внимание!", 32),
    ("", 16),
    ("Auto-ChemInstruct", 24),
    ("Agent-Driven Synthesization of RLHF Data for Chemistry DSLMs", 14),
    ("", 12),
    ("github.com/aayushkrm/auto-cheminstruct", 12),
    ("huggingface.co/datasets/aayushkrm/autochem-instruct", 12),
]:
    p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.alignment = 1

prs.save(OUTPUT)
print(f"Saved {len(prs.slides)} slides → {OUTPUT}")
for i, slide in enumerate(prs.slides):
    t = ""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.text_frame.text.replace("\n", " ")[:70]
            break
    print(f"  {i+1}. {t}")
